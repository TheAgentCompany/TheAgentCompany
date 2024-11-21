import os
from typing import List
from scoring import Result, Checkpoint
from common import evaluate_with_llm, grader
from pptx import Presentation
import logging
import pandas as pd
from collections import defaultdict



@grader
def grade_checkpoint1():
    csv_path = "/workspace/stock_cleaned.csv"
    golden_path = "/utils/stock_data.csv"
    df = pd.read_csv(csv_path)
    golden_df = pd.read_csv(golden_path)
    return df.equals(golden_df)    
    

@grader
def grade_checkpoint2():
    pptx_path = "/workspace/stock_analysis.pptx"
    presentation = Presentation(pptx_path)
    output_dir = "/utils/figs"
    predicates = [
        "Closing Prices Over Time: Display the stock's closing prices over time to show the overall trend.",
        "Close Price with 50 and 200-Day Moving Averages: Overlay the 50-day and 200-day moving averages onto the closing price plot.",
        "Correlation Between Columns: Visualize the correlation matrix focusing on price, volume, and returns.",
        "Daily Return Distribution: Calculate and plot the daily returns to highlight volatility and risk.",
        "Volume Distribution: Plot the distribution of trade volumes to identify trends in stock interest over time.",
    ]

    os.makedirs(output_dir, exist_ok=True)
    
    # key of results are predicates, value of results are the evaluation results
    results = defaultdict(int)

    for idx, slide in enumerate(presentation.slides, start=1):
        image_found = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture type
                image_found = 1
                image = shape.image
                image_bytes = image.blob
                image_path = os.path.join(output_dir, f"slide_{idx}.png")
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
        if image_found:
            predicate = predicates[idx - 1]
            evaluation_result = evaluate_with_llm(None, predicate, image_path=image_path)

            # if there are more than one images on the slide, as long as one image is true, the slide is considered as pass
            if evaluation_result:
                result[predicate] = 1

    idx = 0
    for predicate, result in results:
        if result == 0:
            logging.warning(f"Slide {idx} failed, Predicate: {predicate}")
        idx += 1
    
    return sum([results[pred] for pred in results.keys()])

@grader
def grade_checkpoint3():
    script_path = "/workspace/stock_analysis_script.md"
    predicates = [
        "The script contains clear introductions, addressing and greeting the audience, provides presentation context for the audience, analyzes the business context and impact of each graph. The script ends with closing comments expressing gratitude.",
        "The script includes correspondence for each slide below: Slide 1: Closing Prices Over Time, Slide 2: Close Price with 50 and 200-Day Moving Averages, Slide 3: Correlation Between Columns, Slide 4: Daily Return Distribution, Slide 5: Volume Distribution.",
    ]

    with open(script_path, "r") as f:
        script_content = f.read()

    results = {pred: evaluate_with_llm(script_content, pred) for pred in predicates}

    for pred, result in results.items():
        logging.info(f"{pred}: {'Pass' if result else 'Fail'}")

    # sum over all pass result as int
    return sum([int(result) for result in results.values()])


def grade_checkpoints(trajectory="") -> Result:
    """
    Evaluate all checkpoints and return the results.
    """
    checkpoints: List[Checkpoint] = []
    result = Result(checkpoints)

    checkpoints.append(Checkpoint(1, int(grade_checkpoint1())))
    checkpoints.append(Checkpoint(5, int(grade_checkpoint2())))
    checkpoints.append(Checkpoint(2, int(grade_checkpoint3())))

    return result

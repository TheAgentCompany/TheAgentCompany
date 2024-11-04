import os
import pptx
import requests
from typing import List
from common import create_rocketchat_client, get_chat_history
from scoring import Result, Checkpoint, bonus_for_completing_final

# Initialize RocketChat client
rocket = create_rocketchat_client()


def grade_checkpoint1(trajectory):
    """Check if the agent visits the openhands evaluation page with all supported tasks."""
    return "http://the-agent-company.com:8929/root/openhands/-/tree/main/evaluation?ref_type=heads" in trajectory


def grade_checkpoint2():
    return os.path.exists('/workspace/openhands_intro.pptx')


def get_all_texts(slide):
    """Obtain all text content from the slide."""
    texts = []
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            texts.append(text.lower())
    
    return ' '.join(texts)


def grade_checkpoint3(text: str):
    keywords = [
        "software engineer",
        "swe-bench", "humanevalfix",
        "bird", "biocoder", "ml-bench",
        "apibench", "toolqa", "aiderbench"
    ]
    return all([k in text for k in keywords])

def grade_checkpoint4(text: str):
    keywords = [
        "web browsing",
        "webarena", "miniwob",
    ]
    return all([k in text for k in keywords])

def grade_checkpoint5(text: str):
    keywords = [
        "misc", "assistance",
        "gaia", "gpqa", "agentbench",
        "mint", "eda", "proofwriter"
    ]
    return all([k in text for k in keywords])


def grade_content():
    slide_path = '/workspace/openhands_intro.pptx'
    if not os.path.exists(slide_path):
        return False, False, False
    se, web, misc = "", "", ""
    prs = pptx.Presentation(slide_path)
    for slide in prs.slides:
        slide_texts = get_all_texts(slide)
        if "software engineer" in slide_texts:
            se += slide_texts
        elif "web browsing" in slide_texts:
            web += slide_texts
        elif "misc" in slide_texts:
            misc += slide_texts
    
    passed3 = grade_checkpoint3(se)
    passed4 = grade_checkpoint4(web)
    passed5 = grade_checkpoint5(misc)
    return passed3, passed4, passed5
        
def grade_checkpoint6():
    chat_history = get_chat_history(rocket, "Huang Jie")
    return '/workspace/openhands_intro.pptx' in ' '.join(chat_history)


def grade_checkpoints(trajectory="") -> Result:
    checkpoints: List[Checkpoint] = []
    result = Result(checkpoints, bonus_for_completing_final)

    passed1 = grade_checkpoint1(trajectory)
    checkpoints.append(Checkpoint(1, int(passed1)))

    passed2 = grade_checkpoint2()
    checkpoints.append(Checkpoint(1, int(passed2)))

    passed3, passed4, passed5 = grade_content()
    checkpoints.append(Checkpoint(1, int(passed3)))
    checkpoints.append(Checkpoint(1, int(passed4)))
    checkpoints.append(Checkpoint(1, int(passed5)))

    passed6 = grade_checkpoint6()
    checkpoints.append(Checkpoint(1, int(passed6)))

    return result




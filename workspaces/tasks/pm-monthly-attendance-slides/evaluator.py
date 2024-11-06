import os
import pptx
import requests
from typing import List
from common import create_rocketchat_client
from common import get_rocketchat_personal_chat_history
from scoring import Result, Checkpoint, bonus_for_completing_any
import pandas as pd

# Initialize RocketChat client
rocket = create_rocketchat_client()


def grade_checkpoint1(trajectory):
    """Check if the 'april-attendance-data.xlsx' file on Nextcloud was accessed."""
    required_url_part = "apps/files/files/1233?dir=/Documents/Human%20Resources%20Team/Attendance&openfile=true"
    return required_url_part in trajectory


def grade_checkpoint2():
    """Check if the attendance slide exists at the specified path."""
    return os.path.exists('/workspace/April Attendance.pptx')

def get_all_texts(slide):
    """Obtain all text content from the slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            texts.append(text.lower())
    return ' '.join(texts)

def grade_checkpoint3():
    """Check if the slide contains attendance records for everyone."""
    slide_path = '/workspace/April Attendance.pptx'
    if not os.path.exists(slide_path):
        return False
    
    # Read CSV and get unique names
    df = pd.read_csv('utils/april-attendance-data.csv')
    expected_members = df['Name'].unique().tolist()
    expected_members = [name.lower() for name in expected_members]
    
    # Get slide text content
    prs = pptx.Presentation(slide_path)
    attendance_slide_text = None
    
    # Find the slide titled 'April Attendance'
    for slide in prs.slides:
        slide_text = get_all_texts(slide)
        if 'april attendance' in slide_text.lower():
            attendance_slide_text = slide_text
            break
    
    if not attendance_slide_text:
        return False

    # Check if each member's attendance record is present in the format "name: 10 times"
    return all(f"{member}: 10 times" in attendance_slide_text.lower() for member in expected_members)

def grade_checkpoint4():
    """Check if the agent has a conversation with Chen Xinyi mentioning the correct file path."""
    chat_history = get_rocketchat_personal_chat_history(rocket, "Chen Xinyi")
    return '/workspace/April Attendance.pptx' in ' '.join(chat_history)

def grade_checkpoints(trajectory="") -> Result:
    checkpoints: List[Checkpoint] = []
    result = Result(checkpoints, bonus_for_completing_any)

    passed1 = grade_checkpoint1(trajectory)
    checkpoints.append(Checkpoint(1, int(passed1)))

    passed2 = grade_checkpoint2()
    checkpoints.append(Checkpoint(1, int(passed2)))

    passed3 = grade_checkpoint3()
    checkpoints.append(Checkpoint(1, int(passed3)))

    passed4 = grade_checkpoint4()
    checkpoints.append(Checkpoint(1, int(passed4)))

    return result




import os
from typing import List
from config import *
from scoring import Result, Checkpoint, bonus_for_completing_final
import pptx

expected_repos = [
    {"name": "api-server", "description": "", "issues": 0, "merge_requests": 0},
    {"name": "bustub", "description": "The BusTub Relational Database Management System (Educational)", "issues": 28, "merge_requests": 14},
    {"name": "colly", "description": "Elegant Scraper and Crawler Framework for Golang", "issues": 152, "merge_requests": 41},
    {"name": "copilot-arena-server", "description": "", "issues": 0, "merge_requests": 0},
    {"name": "Documentation", "description": "Wiki for company-wide doc", "issues": 0, "merge_requests": 0},
    {"name": "janusgraph", "description": "JanusGraph: an open-source, distributed graph database", "issues": 491, "merge_requests": 40},
    {"name": "llama.cpp", "description": "LLM inference in C/C++", "issues": 268, "merge_requests": 277},
    {"name": "node-red", "description": "Low-code programming for event-driven applications", "issues": 304, "merge_requests": 95},
    {"name": "openhands", "description": "OpenHands: Code Less, Make More", "issues": 120, "merge_requests": 34},
    {"name": "opensearch", "description": "Open source distributed and RESTful search engine.", "issues": 1700, "merge_requests": 161},
    {"name": "raft", "description": "", "issues": 0, "merge_requests": 0},
    {"name": "risingwave", "description": "Best-in-class stream processing, analytics, and management. Unified streaming and batch. PostgreSQL compatible.", "issues": 996, "merge_requests": 88},
    {"name": "sotopia", "description": "Sotopia: an Open-ended Social Learning Environment (ICLR 2024 spotlight)", "issues": 12, "merge_requests": 6},
    {"name": "streamlit", "description": "Streamlit — A faster way to build and share data apps.", "issues": 919, "merge_requests": 38},
]


def get_all_texts(slide):
    """Extract all text from a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text.lower())
    return ' '.join(texts)

def grade_checkpoint1():
    """Check if the presentation file exists."""
    return os.path.exists('/workspace/Engineering_Progress.pptx')

def grade_checkpoint2():
    """Check each repository slide for the correct name, description, issues, and merge requests."""
    slide_path = '/workspace/Engineering_Progress.pptx'
    if not os.path.exists(slide_path):
        return 0  # No presentation file found, so no points
    
    prs = pptx.Presentation(slide_path)
    total_points = 0
    #Checkpoint 2 gives no partial points per component i.e. all 4 have to be correctly present
    for repo in expected_repos:
        repo_correct = False
        for slide in prs.slides:
            slide_text = get_all_texts(slide)
            # Verify all four components: name, description, issues, merge requests
            if (repo["name"].lower() in slide_text and
                repo["description"].lower() in slide_text and
                f"issues: {repo['issues']}" in slide_text and
                f"merge requests: {repo['merge_requests']}" in slide_text):
                repo_correct = True
                break  # Found correct slide for this repo, move to the next repo
        
        if repo_correct:
            total_points += 1  # Award 1 point if all components are correct

    return total_points


def grade_checkpoints(trajectory=""):
    checkpoints: List[Checkpoint] = []
    result = Result(checkpoints, bonus_for_completing_final)

    # Checkpoint 1: Check if the presentation file exists
    passed1 = grade_checkpoint1()
    checkpoints.append(Checkpoint(1, int(passed1)))

    # Checkpoint 2: Verify each repository slide has correct details
    passed2 = grade_checkpoint2()
    checkpoints.append(Checkpoint(14, passed2))

    return result


import os
import logging
import pptx
from typing import List

from common import get_all_texts_from_slide, evaluate_chat_history_with_llm, create_rocketchat_client, grader
from scoring import Result, Checkpoint, bonus_for_completing_final

SLIDES_PATH = "/workspace/engineers.pptx"
NPC_NAME = "chen_xinyi"

ENGINEERS = [
    ["Sarah", "Johnson", "CTO"],
    ["Li", "Ming", "Database Team Project Manager"],
    ["Zhang", "Wei", "Senior Software Engineer"],
    ["Wang", "Fang", "AI Researcher"],
    ["Mike", "Chen", "Senior Software Engineer"],
    ["Emily", "Zhou", "Software Engineer"]    
]

@grader
def grade_checkpoint1():
    return os.path.exists(SLIDES_PATH)

@grader
def grade_checkpoint2(slides):
    """Slides should contain information about all engineers"""
    if slides is None or len(slides.slides) < 1:
        return False
    text = ""
    for i in range(len(slides.slides)):
        try:
            text += get_all_texts_from_slide(slides.slides[i])
        except Exception as e:
            logging.error("Error reading slide 2: %s", e)
            return False
    for engineer in ENGINEERS:
        if not all(k.lower() in text for k in engineer):
            return False
    return True

@grader
def grade_checkpoint3():
    rocket_client = create_rocketchat_client()
    return evaluate_chat_history_with_llm(rocket_client, NPC_NAME, "introductory slides are finished")


def grade_checkpoints(trajectory="") -> Result:
    checkpoint1 = Checkpoint(1, 1*int(grade_checkpoint1()))
    try:
        slides = pptx.Presentation(SLIDES_PATH)
    except Exception as e:
        logging.error("Error reading slides: %s", e)
        slides = None
    checkpoint2 = Checkpoint(2, 2*int(grade_checkpoint2(slides)))
    checkpoint3 = Checkpoint(1, 1*int(grade_checkpoint3()))

    return Result([checkpoint1, checkpoint2, checkpoint3])

if __name__ == "__main__":
    print(grade_checkpoints().to_dict())